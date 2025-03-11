import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
import PyPDF2
from docx import Document
import textract

app = Flask(__name__)
CORS(app)  # Allow CORS for all routes

# Function to extract text from PowerPoint files
def extract_text_from_ppt(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error processing PPT file: {str(e)}"

# Function to extract text from PDF files
def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = []
            for page in pdf_reader.pages:
                text.append(page.extract_text())
            return "\n".join(text)
    except Exception as e:
        return f"Error processing PDF file: {str(e)}"

# Function to extract text from DOCX files
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error processing DOCX file: {str(e)}"

# Function to extract text from DOC files using textract
def extract_text_from_doc(file_path):
    try:
        text = textract.process(file_path).decode('utf-8')
        return text
    except Exception as e:
        return f"Error processing DOC file: {str(e)}"

# Function to determine file type and extract text
def extract_text_from_file(file_path):
    file_extension = os.path.splitext(file_path.lower())[1]
    
    if file_extension == '.pptx':
        return extract_text_from_ppt(file_path)
    elif file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.doc':
        return extract_text_from_doc(file_path)
    else:
        return f"Unsupported file format: {file_extension}"

@app.route('/extract-text', methods=['POST'])
def extract_text():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    file_path = os.path.join("uploads", file.filename)
    os.makedirs("uploads", exist_ok=True)
    file.save(file_path)
    
    extracted_text = extract_text_from_file(file_path)
    os.remove(file_path)  # Cleanup after processing
    
    return jsonify({"filename": file.filename, "text": extracted_text})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5002, debug=True)
